from __future__ import annotations
from dataclasses import dataclass
from pathlib import Path
from pptx import Presentation


@dataclass
class FontInfo:
    major: str = "Calibri"
    minor: str = "Calibri"


@dataclass
class ColorScheme:
    accent1: str = "4472C4"
    accent2: str = "ED7D31"
    accent3: str = "70AD47"
    accent4: str = "FFC000"
    dark1: str = "000000"
    dark2: str = "44546A"
    light1: str = "FFFFFF"
    light2: str = "F2F2F2"


@dataclass
class TemplateInfo:
    slide_width: int
    slide_height: int
    color_scheme: ColorScheme
    font: FontInfo
    logo_path: str | None = None
    background_style: str = "solid"


def analyze_template(path: str) -> TemplateInfo:
    file_path = Path(path)
    if not file_path.exists():
        raise FileNotFoundError(f"Template not found: {path}")
    prs = Presentation(path)
    width = prs.slide_width
    height = prs.slide_height
    scheme = ColorScheme()
    font = FontInfo()
    for master in prs.slide_masters:
        if hasattr(master, 'slide_layouts'):
            pass
    logo_path = None
    if prs.slides:
        for shape in prs.slides[0].shapes:
            if shape.shape_type == 13:
                logo_path = f"extracted_logo_{shape.shape_id}.png"
    return TemplateInfo(
        slide_width=width,
        slide_height=height,
        color_scheme=scheme,
        font=font,
        logo_path=logo_path,
    )

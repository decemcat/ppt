from __future__ import annotations
from collections import Counter
from pptx import Presentation
from pptx.util import Pt
from ppt_agent.style.profile import StyleProfile, ColorScheme, FontProfile, LayoutRatios


def _rgb_to_hex(rgb) -> str:
    if rgb is None:
        return "#333333"
    return f"#{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"


def _most_common(items: list) -> str:
    if not items:
        return "#333333"
    counter = Counter(items)
    return counter.most_common(1)[0][0]


class StyleExtractor:
    @staticmethod
    def extract(pptx_path: str, name: str = "extracted") -> StyleProfile:
        prs = Presentation(pptx_path)
        colors = StyleExtractor._extract_colors(prs)
        fonts = StyleExtractor._extract_fonts(prs)
        layout = StyleExtractor._extract_layout(prs)
        return StyleProfile(
            name=name,
            colors=colors,
            fonts=fonts,
            layout=layout,
            source_file=pptx_path,
        )

    @staticmethod
    def _extract_colors(prs: Presentation) -> ColorScheme:
        bg_colors = []
        text_colors = []
        fill_colors = []
        for slide in prs.slides:
            try:
                if slide.background.fill.fore_color.rgb:
                    bg_colors.append(_rgb_to_hex(slide.background.fill.fore_color.rgb))
            except Exception:
                pass
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            try:
                                if run.font.color and run.font.color.rgb:
                                    text_colors.append(_rgb_to_hex(run.font.color.rgb))
                            except Exception:
                                pass
                try:
                    if hasattr(shape, "fill") and shape.fill.fore_color and shape.fill.fore_color.rgb:
                        fill_colors.append(_rgb_to_hex(shape.fill.fore_color.rgb))
                except Exception:
                    pass
        background = _most_common(bg_colors) if bg_colors else "#FFFFFF"
        primary = _most_common(fill_colors) if fill_colors else "#1F4E79"
        text_primary = _most_common(text_colors) if text_colors else "#333333"
        secondary = fill_colors[1] if len(fill_colors) > 1 else primary
        accent = "#D6E4F0"
        try:
            if prs.slide_masters[0].background.fill.fore_color.rgb:
                accent = _rgb_to_hex(prs.slide_masters[0].background.fill.fore_color.rgb)
        except Exception:
            pass
        return ColorScheme(
            primary=primary,
            secondary=secondary,
            accent=accent,
            background=background,
            text_primary=text_primary,
            text_secondary="#333333",
            text_light="#FFFFFF",
        )

    @staticmethod
    def _extract_fonts(prs: Presentation) -> FontProfile:
        title_fonts = []
        title_sizes = []
        body_fonts = []
        body_sizes = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        font_name = run.font.name
                        font_size = run.font.size
                        if font_name:
                            if font_size and font_size >= Pt(24):
                                title_fonts.append(font_name)
                                title_sizes.append(int(font_size / Pt(1)))
                            else:
                                body_fonts.append(font_name)
                                if font_size:
                                    body_sizes.append(int(font_size / Pt(1)))
        title_font = _most_common(title_fonts) if title_fonts else "微软雅黑"
        body_font = _most_common(body_fonts) if body_fonts else "微软雅黑"
        title_size = _most_common(title_sizes) if title_sizes else 32
        body_size_val = _most_common(body_sizes) if body_sizes else 14
        return FontProfile(
            title_font=title_font,
            title_size=title_size,
            subtitle_font=title_font,
            subtitle_size=max(int(title_size * 0.65), 14),
            body_font=body_font,
            body_size=body_size_val,
            caption_font=body_font,
            caption_size=max(int(body_size_val * 0.8), 10),
        )

    @staticmethod
    def _extract_layout(prs: Presentation) -> LayoutRatios:
        from pptx.util import Emu
        slide_height = prs.slide_height
        title_heights = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame and shape.top < Emu(slide_height // 3):
                    max_font = 0
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.size and run.font.size > max_font:
                                max_font = run.font.size
                    if max_font >= Pt(24):
                        title_heights.append(shape.top + shape.height)
        if title_heights:
            avg_title_h = sum(title_heights) / len(title_heights)
            title_ratio = round(float(avg_title_h / slide_height), 2)
        else:
            title_ratio = 0.15
        return LayoutRatios(
            title_height_ratio=title_ratio,
            content_margin_ratio=0.08,
            column_gap_ratio=0.05,
            element_spacing_ratio=0.03,
        )

from __future__ import annotations
import tempfile
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from typing import TYPE_CHECKING

from ppt_agent.models import SlideContent, SlideFramework, PPTFramework
from ppt_agent.template.analyzer import TemplateInfo, analyze_template
from ppt_agent.generator.shape_renderer import render_diagram_to_slide
from ppt_agent.generator.image_fallback import render_as_image
from ppt_agent.generator.quality_check import score_quality

if TYPE_CHECKING:
    from ppt_agent.style.profile import StyleProfile
    from ppt_agent.generator.image_gen import ImageGenerator

IMAGE_FALLBACK_THRESHOLD = 60
MAX_AI_IMAGES = 2


def _hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def generate_pptx(
    ppt_framework: PPTFramework,
    template_path: str,
    output_path: str = "output.pptx",
    style_profile: StyleProfile | None = None,
    image_gen: ImageGenerator | None = None,
) -> str:
    """Generate a .pptx file from a framework.

    Returns the output path.
    """
    template_info = analyze_template(template_path)
    prs = Presentation(template_path)
    blank_layout = prs.slide_layouts[0]
    ai_counter = {"ai_count": 0}
    for slide_content in ppt_framework.framework.slides:
        slide = prs.slides.add_slide(blank_layout)
        if slide_content.slide_type == "title":
            _render_title_slide(slide, slide_content, template_info, style_profile)
        elif slide_content.slide_type == "section_header":
            _render_section_header(slide, slide_content, template_info, style_profile)
        elif slide_content.slide_type == "arch_diagram" and slide_content.diagram:
            _render_arch_slide(slide, slide_content, template_info, style_profile, image_gen, ai_counter)
        else:
            _render_text_slide(slide, slide_content, template_info, style_profile, image_gen, ai_counter)
    prs.save(output_path)
    return output_path


def _render_slide_title(slide, title: str, template: TemplateInfo, size: int = 24, style_profile: StyleProfile | None = None):
    for shape in slide.shapes:
        if hasattr(shape, "text_frame"):
            p = shape.text_frame.paragraphs[0]
            p.text = title
            if style_profile:
                p.font.size = Pt(style_profile.fonts.title_size)
                p.font.name = style_profile.fonts.title_font
                p.font.color.rgb = _hex_to_rgb(style_profile.colors.text_primary.lstrip("#") if style_profile.colors.text_primary.startswith("#") else style_profile.colors.text_primary)
            else:
                p.font.size = Pt(size)
                p.font.color.rgb = _hex_to_rgb(template.color_scheme.dark1)
            break


def _render_title_slide(slide, content: SlideContent, template: TemplateInfo, style_profile: StyleProfile | None = None):
    _render_slide_title(slide, content.title, template, size=28, style_profile=style_profile)
    if content.notes:
        slide.notes_slide.notes_text_frame.text = content.notes


def _render_section_header(slide, content: SlideContent, template: TemplateInfo, style_profile: StyleProfile | None = None):
    _render_slide_title(slide, content.title, template, style_profile=style_profile)
    if content.notes:
        slide.notes_slide.notes_text_frame.text = content.notes


def _render_text_slide(slide, content: SlideContent, template: TemplateInfo, style_profile: StyleProfile | None = None, image_gen: ImageGenerator | None = None, ai_counter: dict[str, int] | None = None):
    _render_slide_title(slide, content.title, template, size=24, style_profile=style_profile)
    if content.bullets:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.has_text_frame:
                tf = shape.text_frame
                tf.clear()
                for i, bullet in enumerate(content.bullets):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = f"• {bullet}"
                    if style_profile:
                        p.font.size = Pt(style_profile.fonts.body_size)
                        p.font.name = style_profile.fonts.body_font
                        p.font.color.rgb = _hex_to_rgb(style_profile.colors.text_secondary.lstrip("#") if style_profile.colors.text_secondary.startswith("#") else style_profile.colors.text_secondary)
                    else:
                        p.font.size = Pt(16)
                        p.font.color.rgb = _hex_to_rgb(template.color_scheme.dark2)
                break
    if ai_counter is not None:
        _maybe_insert_ai_image(slide, content, image_gen, ai_counter)
    if content.notes:
        slide.notes_slide.notes_text_frame.text = content.notes


def _render_arch_slide(slide, content: SlideContent, template: TemplateInfo, style_profile: StyleProfile | None = None, image_gen: ImageGenerator | None = None, ai_counter: dict[str, int] | None = None):
    _render_slide_title(slide, content.title, template, style_profile=style_profile)
    diag = content.diagram
    if image_gen and content.image_prompt and ai_counter is not None:
        _maybe_insert_ai_image(slide, content, image_gen, ai_counter)
    else:
        render_diagram_to_slide(slide, diag, template, style_profile=style_profile)
        quality = score_quality(slide, template.slide_width, template.slide_height)
        if quality < IMAGE_FALLBACK_THRESHOLD:
            _clear_slide_shapes(slide)
            _render_slide_title(slide, content.title, template, size=24, style_profile=style_profile)
            tmp_dir = Path(tempfile.mkdtemp())
            img_path = str(tmp_dir / "arch.png")
            render_as_image(diag, img_path)
            if Path(img_path).exists():
                slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), width=Inches(11))
    if content.notes:
        slide.notes_slide.notes_text_frame.text = content.notes


def _maybe_insert_ai_image(slide, content: SlideContent, image_gen: ImageGenerator | None, counter: dict[str, int]):
    if not image_gen or not content.image_prompt:
        return
    if counter["ai_count"] >= MAX_AI_IMAGES:
        return
    tmp_dir = tempfile.mkdtemp()
    img_path = str(Path(tmp_dir) / "ai_gen.png")
    result = image_gen.generate(content.image_prompt, size="1024x1024", output_path=img_path)
    if result and Path(result).exists():
        slide.shapes.add_picture(result, Inches(7), Inches(1.5), width=Inches(5), height=Inches(5))
        counter["ai_count"] += 1


def _clear_slide_shapes(slide):
    for shape in list(slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

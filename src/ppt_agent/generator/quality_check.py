from pptx import Presentation
from pptx.util import Emu


def _shape_bounds(shape):
    """Get (left, top, right, bottom) in EMU."""
    return (
        shape.left,
        shape.top,
        shape.left + shape.width,
        shape.top + shape.height,
    )


def _rects_overlap(a, b):
    """Check if two axis-aligned rectangles overlap."""
    return not (a[2] <= b[0] or b[2] <= a[0] or a[3] <= b[1] or b[3] <= a[1])


def check_overlap(slide) -> list[tuple]:
    """Return list of overlapping shape index pairs."""
    shapes = list(slide.shapes)
    overlaps = []
    rects = [_shape_bounds(s) for s in shapes]
    for i in range(len(shapes)):
        for j in range(i + 1, len(shapes)):
            if _rects_overlap(rects[i], rects[j]):
                overlaps.append((i, j))
    return overlaps


def check_slide_bounds(slide, slide_width: int, slide_height: int) -> list[str]:
    """Return shapes that exceed slide boundaries."""
    violations = []
    for shape in slide.shapes:
        if (shape.left < 0 or shape.top < 0 or
                shape.left + shape.width > slide_width or
                shape.top + shape.height > slide_height):
            violations.append(str(shape))
    return violations


def check_font_size(slide, min_pt: int = 8) -> list[str]:
    """Return shapes with font size below min_pt."""
    violations = []
    for shape in slide.shapes:
        if hasattr(shape, "text_frame"):
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size and run.font.size < min_pt * 12700:
                        violations.append(str(run))
    return violations


def score_quality(slide, slide_width: int, slide_height: int) -> int:
    """Score 0-100. Below 60 triggers Image Fallback."""
    score = 100
    overlaps = check_overlap(slide)
    if overlaps:
        score -= min(40, len(overlaps) * 20)
    bounds_violations = check_slide_bounds(slide, slide_width, slide_height)
    if bounds_violations:
        score -= min(30, len(bounds_violations) * 15)
    font_violations = check_font_size(slide)
    if font_violations:
        score -= min(20, len(font_violations) * 5)
    return max(score, 0)

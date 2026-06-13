from __future__ import annotations
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_CONNECTOR_TYPE
from pptx.oxml.ns import qn

from ppt_agent.models import ArchDiagram, ArchNode, ArchEdge
from ppt_agent.template.analyzer import TemplateInfo


def _hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def layout_layered(diag: ArchDiagram) -> dict[str, dict]:
    """Top-to-bottom layered layout for arch diagrams."""
    positions = {}
    group_nodes = {n.id: n for n in diag.nodes if n.style == "group"}
    component_nodes = {n.id: n for n in diag.nodes if n.style != "group"}
    y_offset = Emu(500000)
    for gid, gnode in group_nodes.items():
        positions[gid] = {"x": Emu(1000000), "y": y_offset, "w": Emu(11000000), "h": Emu(600000)}
        y_offset += Emu(900000)
    for cid in component_nodes:
        positions[cid] = {"x": Emu(1000000), "y": y_offset, "w": Emu(11000000), "h": Emu(500000)}
        y_offset += Emu(800000)
    for node in diag.nodes:
        if node.id not in positions:
            positions[node.id] = {"x": Emu(1000000), "y": y_offset, "w": Emu(2500000), "h": Emu(500000)}
            y_offset += Emu(700000)
    return positions


def layout_flow(diag: ArchDiagram) -> dict[str, dict]:
    """Left-to-right flow layout."""
    positions = {}
    x_offset = Emu(500000)
    for node in diag.nodes:
        positions[node.id] = {
            "x": x_offset,
            "y": Emu(2000000),
            "w": Emu(2500000),
            "h": Emu(1000000),
        }
        x_offset += Emu(3000000)
    return positions


def layout_radial(diag: ArchDiagram) -> dict[str, dict]:
    """Simple radial layout (center + surrounding)."""
    import math
    positions = {}
    cx = Emu(6500000)
    cy = Emu(3500000)
    if diag.nodes:
        center = diag.nodes[0]
        positions[center.id] = {"x": cx - Emu(1000000), "y": cy - Emu(350000), "w": Emu(2000000), "h": Emu(700000)}
        angle_step = 360 / max(len(diag.nodes) - 1, 1)
        for i, node in enumerate(diag.nodes[1:], 1):
            rad = math.radians(i * angle_step)
            rx = int(cx + Emu(2500000) * math.cos(rad) - Emu(1000000))
            ry = int(cy + Emu(2000000) * math.sin(rad) - Emu(300000))
            positions[node.id] = {"x": rx, "y": ry, "w": Emu(2000000), "h": Emu(600000)}
    return positions


def layout_grid(diag: ArchDiagram) -> dict[str, dict]:
    """Grid layout for equal-sized boxes."""
    positions = {}
    cols = 3
    cell_w = Emu(3800000)
    cell_h = Emu(1500000)
    start_x = Emu(500000)
    start_y = Emu(1000000)
    for i, node in enumerate(diag.nodes):
        col = i % cols
        row = i // cols
        positions[node.id] = {
            "x": start_x + col * cell_w,
            "y": start_y + row * cell_h,
            "w": cell_w - Emu(200000),
            "h": cell_h - Emu(200000),
        }
    return positions


def _get_layout_func(diagram_type: str):
    return {
        "layered": layout_layered,
        "flow": layout_flow,
        "radial": layout_radial,
        "grid": layout_grid,
    }.get(diagram_type, layout_layered)


def render_diagram_to_slide(slide, diag: ArchDiagram, template: TemplateInfo):
    """Render an ArchDiagram onto a slide using python-pptx shapes."""
    layout_func = _get_layout_func(diag.type)
    positions = layout_func(diag)
    accent = _hex_to_rgb(template.color_scheme.accent1)
    dark = _hex_to_rgb(template.color_scheme.dark1)
    light = _hex_to_rgb(template.color_scheme.light1)
    node_map = {n.id: n for n in diag.nodes}
    for node in diag.nodes:
        pos = positions.get(node.id)
        if not pos:
            continue
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            pos["x"], pos["y"], pos["w"], pos["h"],
        )
        shape.fill.solid()
        if node.style == "group":
            shape.fill.fore_color.rgb = accent
            shape.line.fill.background()
        else:
            shape.fill.fore_color.rgb = light
            shape.line.color.rgb = accent
            shape.line.width = Pt(1)
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = node.label
        p.font.size = Pt(12)
        p.font.color.rgb = light if node.style == "group" else dark
        p.alignment = PP_ALIGN.CENTER
    for edge in diag.edges:
        from_pos = positions.get(edge.from_id)
        to_pos = positions.get(edge.to_id)
        if not from_pos or not to_pos:
            continue
        from_cx = from_pos["x"] + from_pos["w"] // 2
        from_cy = from_pos["y"] + from_pos["h"]
        to_cx = to_pos["x"] + to_pos["w"] // 2
        to_cy = to_pos["y"]
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR_TYPE.STRAIGHT,
            from_cx, from_cy, to_cx, to_cy,
        )
        connector.line.color.rgb = accent
        connector.line.width = Pt(1.5)
        if edge.label:
            # Add a text label near the midpoint
            mid_x = (from_cx + to_cx) // 2
            mid_y = (from_cy + to_cy) // 2
            label_box = slide.shapes.add_textbox(mid_x, mid_y - Emu(100000), Emu(800000), Emu(200000))
            tf = label_box.text_frame
            p = tf.paragraphs[0]
            p.text = edge.label
            p.font.size = Pt(9)
            p.font.color.rgb = accent

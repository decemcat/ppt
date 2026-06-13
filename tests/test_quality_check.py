from ppt_agent.generator.quality_check import (
    check_overlap, check_slide_bounds, check_font_size, score_quality
)
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


class TestQualityCheck:
    def test_no_overlap(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        left = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
        right = slide.shapes.add_textbox(Inches(5), Inches(1), Inches(3), Inches(1))
        overlaps = check_overlap(slide)
        assert len(overlaps) == 0

    def test_has_overlap(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        a = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
        b = slide.shapes.add_textbox(Inches(2), Inches(1.5), Inches(4), Inches(2))
        overlaps = check_overlap(slide)
        assert len(overlaps) > 0

    def test_within_bounds(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        s = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(5), Inches(3))
        violations = check_slide_bounds(slide, prs.slide_width, prs.slide_height)
        assert len(violations) == 0

    def test_out_of_bounds(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        s = slide.shapes.add_textbox(Inches(12), Inches(0), Inches(5), Inches(3))
        violations = check_slide_bounds(slide, prs.slide_width, prs.slide_height)
        assert len(violations) > 0

    def test_score_perfect(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
        score = score_quality(slide, prs.slide_width, prs.slide_height)
        assert score == 100

    def test_score_with_overlap(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        a = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
        b = slide.shapes.add_textbox(Inches(2), Inches(1.5), Inches(4), Inches(2))
        score = score_quality(slide, prs.slide_width, prs.slide_height)
        assert score < 100

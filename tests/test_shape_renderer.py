from ppt_agent.models import ArchDiagram, ArchNode, ArchEdge
from ppt_agent.generator.shape_renderer import (
    layout_layered, layout_flow, layout_radial, layout_grid,
    render_diagram_to_slide,
)
from ppt_agent.template.analyzer import TemplateInfo, ColorScheme, FontInfo
from pptx import Presentation
from pptx.util import Inches


class TestShapeRenderer:
    def test_layout_layered(self):
        diag = ArchDiagram(
            type="layered",
            nodes=[
                ArchNode(id="app", label="应用层", children=["web", "api"]),
                ArchNode(id="data", label="数据层"),
            ],
            edges=[ArchEdge(from_id="app", to_id="data")],
        )
        positions = layout_layered(diag)
        assert "app" in positions
        assert positions["app"]["y"] < positions["data"]["y"]

    def test_render_to_slide(self):
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        diag = ArchDiagram(
            type="layered",
            nodes=[ArchNode(id="a", label="A"), ArchNode(id="b", label="B")],
            edges=[ArchEdge(from_id="a", to_id="b")],
        )
        template = TemplateInfo(
            slide_width=prs.slide_width,
            slide_height=prs.slide_height,
            color_scheme=ColorScheme(),
            font=FontInfo(),
        )
        render_diagram_to_slide(slide, diag, template)
        assert len(slide.shapes) > 0
